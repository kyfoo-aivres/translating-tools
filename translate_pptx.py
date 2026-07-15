import argparse
import os
import shutil
import glob
from pptx import Presentation
from googletrans import Translator
from tqdm import tqdm
import asyncio

def extract_paragraphs(shapes):
    """
    Recursively extract all paragraphs from presentation shapes.
    This handles normal shapes, grouped shapes, and tables.
    """
    paragraphs = []
    for shape in shapes:
        # Check if the shape is a group that contains other shapes
        if hasattr(shape, "shapes"):
            paragraphs.extend(extract_paragraphs(shape.shapes))
            
        # Check if the shape has a text frame
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraphs.append(paragraph)
                        
        # Check if the shape contains a table
        if hasattr(shape, "has_table") and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame"):
                        for paragraph in cell.text_frame.paragraphs:
                            paragraphs.append(paragraph)
    return paragraphs

def translate_presentation(input_file, src_lang="zh-cn", dest_lang="en"):
    # Ensure input file exists
    if not os.path.exists(input_file):
        print(f"Error: The file '{input_file}' was not found.")
        return

    # Prepare input, output, and log paths
    file_directory = os.path.dirname(os.path.abspath(input_file))
    base_name = os.path.splitext(os.path.basename(input_file))[0]
    ext_name = os.path.splitext(input_file)[1]
    
    # Setup paths according to the requirements
    output_file = os.path.join(file_directory, f"{base_name}_{dest_lang}{ext_name}")
    issue_file = os.path.join(file_directory, f"Issues_{base_name}.txt")

    # Duplicate the source PPTX
    try:
        shutil.copy2(input_file, output_file)
        print(f"Copied source file to: {output_file}")
    except Exception as e:
        print(f"Error copying the file: {e}")
        return

    # Load the duplicated PPTX
    try:
        prs = Presentation(output_file)
    except Exception as e:
        print(f"Error opening copied presentation (make sure it's a valid .pptx): {e}")
        return

    # Translator will be initialized inside the async loop if needed
    # Extract all textual components that need to be translated
    print("Analyzing presentation and extracting texts...")
    all_runs = []
    notes_runs = []
    all_paragraphs = []
    
    for slide in prs.slides:
        paragraphs = extract_paragraphs(slide.shapes)
        all_paragraphs.extend(paragraphs)
        
        for paragraph in paragraphs:
            for run in paragraph.runs:
                if run.text and run.text.strip():
                    all_runs.append(run)
                    
        if slide.has_notes_slide:
            notes_text_frame = slide.notes_slide.notes_text_frame
            for paragraph in notes_text_frame.paragraphs:
                all_paragraphs.append(paragraph)
                for run in paragraph.runs:
                    if run.text and run.text.strip():
                        notes_runs.append(run)
        
    print(f"Discovered {len(all_runs)} text runs in slides to translate.")
    if notes_runs:
        print(f"Discovered {len(notes_runs)} text runs in notes to translate.")
    
    # Setup issue logging 
    if os.path.exists(issue_file):
        os.remove(issue_file) # Clean up old issues file if it exists
        
    issue_found = False

    
    
    async def translate_run_list(runs_list, desc):
        nonlocal issue_found
        # Initialize the translator inside the event loop
        translator = Translator()
        
        # Proceed with the translation and show the progress bar using tqdm
        for run in tqdm(runs_list, desc=desc, unit="run"):
            original_text = run.text
            max_retries = 3
            retry_count = 0
            
            while retry_count <= max_retries:
                try:
                    # Perform the translation
                    translation = translator.translate(original_text, src=src_lang, dest=dest_lang)
                    
                    # Wait for the result if it's an async coroutine
                    if asyncio.iscoroutine(translation):
                        translation = await translation
                        
                    # Update the run's text with the translated text
                    if translation and hasattr(translation, 'text'):
                        run.text = translation.text
                        
                    break # Success, exit retry loop
                except Exception as e:
                    error_msg = str(e)
                    if ("[Errno 11001] getaddrinfo failed" in error_msg or not error_msg.strip()) and retry_count < max_retries:
                        retry_count += 1
                        await asyncio.sleep(5)
                        translator = Translator() # Re-initialize translator on error just in case
                        continue
                    
                    issue_found = True
                    # Log the issue into the issue file
                    with open(issue_file, "a", encoding="utf-8") as file:
                        file.write(f"Failed to translate block ({desc}): '{original_text}'\n")
                        if retry_count > 0:
                            file.write(f"Failed after {retry_count} retries.\n")
                        file.write(f"Error details: {error_msg}\n")
                        file.write("-" * 50 + "\n")
                    break
                    
    async def main_translation_task():
        if all_runs:
            await translate_run_list(all_runs, "Translating Slides")
        if notes_runs:
            await translate_run_list(notes_runs, "Translating Notes")

    # Execute the async translation process
    asyncio.run(main_translation_task())
    
    # Post-translation: Fix missing spaces between runs in Latin-based translations
    cjk_langs = ['zh', 'zh-cn', 'zh-tw', 'ja', 'ko']
    if dest_lang.lower() not in cjk_langs:
        no_space_before = tuple('.,!?:;)]}”’')
        no_space_after = tuple('([{"‘')
        
        for paragraph in all_paragraphs:
            runs = paragraph.runs
            if not runs:
                continue
                
            for i in range(len(runs) - 1):
                curr_run = runs[i]
                next_run = runs[i+1]
                
                c_text = curr_run.text
                n_text = next_run.text
                
                if not c_text or not n_text:
                    continue
                    
                # If current run doesn't end with a space, and next doesn't start with a space
                if not c_text[-1].isspace() and not n_text[0].isspace():
                    # Add a space to current_run unless it breaks punctuation spacing rules
                    if not n_text.startswith(no_space_before) and not c_text.endswith(no_space_after):
                        curr_run.text = c_text + " "
                
    # Save the updated presentation
    try:
        prs.save(output_file)
        print(f"\nTranslation complete. The translated file is ready at: {output_file}")
    except Exception as e:
        print(f"\nError: Could not save the translated presentation: {e}")

    # Notify the user if issues were logged
    if issue_found:
        print(f"Note: Some issues occurred during translation. See '{issue_file}' for details.")

def process_file(file_path, src_lang, dest_lang):
    print(f"\n{'='*50}\nProcessing file: {file_path}\n{'='*50}")
    translate_presentation(file_path, src_lang, dest_lang)

def main():
    parser = argparse.ArgumentParser(description="Translate Microsoft PowerPoint (.pptx) files.")
    
    # Mutually exclusive group for file or directory input
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("-f", "--file", help="Path to a single source PPTX file to translate")
    group.add_argument("-d", "--directory", help="Path to a directory containing PPTX files to translate")
    
    parser.add_argument("--source", default="zh-cn", help="Source language code (default: zh-CN)")
    parser.add_argument("--target", default="en", help="Target language code (default: en)")
    
    args = parser.parse_args()
    
    if args.file:
        process_file(args.file, args.source, args.target)
    elif args.directory:
        if not os.path.isdir(args.directory):
            print(f"Error: Directory '{args.directory}' does not exist.")
            return
            
        pptx_files = glob.glob(os.path.join(args.directory, "*.pptx"))
        
        if not pptx_files:
            print(f"No .pptx files found in directory: {args.directory}")
            return
            
        print(f"Discovered {len(pptx_files)} .pptx files in '{args.directory}'.")
        for pptx in pptx_files:
            # Skip hidden/temporary PowerPoint files
            if not os.path.basename(pptx).startswith("~$"):
                process_file(pptx, args.source, args.target)

if __name__ == "__main__":
    main()
